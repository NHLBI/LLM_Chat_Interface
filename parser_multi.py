#!/usr/bin/python3
"""
NHLBI Chat – multipurpose document parser with optional on-prem OCR support
==========================================================================
• Text / markdown / json / xml  (.txt .md .json .xml)
• Word  (.docx)        – text, tables, images→OCR
• PowerPoint (.pptx)    – text frames, tables, images→OCR
• PDF   (.pdf)          – page text plus images→OCR
• CSV / Excel (.csv .xls .xlsx) – unchanged (json serialization)

Set environment variables:
  OCR_ENABLED=0            # disables all OCR work
  OCR_LANG=spa+eng         # any language string valid for tesseract
"""

import os, sys, io, logging, itertools
import pandas as pd

# 3rd-party helpers ----------------------------------------------------
from PIL import Image
import pytesseract
import fitz  # PyMuPDF
from docx import Document
from docx.text.paragraph import Paragraph
from docx.table import Table
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import xlrd

# ---------- configuration ----------
OCR_ENABLED = os.getenv("OCR_ENABLED", "1") != "0"
OCR_LANG = os.getenv("OCR_LANG", "eng")

log = logging.getLogger("nhlbi_parser")
log.setLevel(logging.INFO)


def ocr_image_bytes(image_bytes: bytes, description: str = "") -> str:
    """
    Run Tesseract OCR on raw image bytes. Returns empty string if OCR disabled
    or no text found.
    """
    if not OCR_ENABLED:
        return ""
    try:
        img = Image.open(io.BytesIO(image_bytes))
        text = pytesseract.image_to_string(img, lang=OCR_LANG)
        return text.strip()
    except Exception as exc:  # pylint: disable=broad-except
        log.warning("OCR failure on %s: %s", description, exc)
        return ""


# ==========================================================
# Dispatch entry
# ==========================================================
def parse_doc(file_path: str, filename: str) -> str:
    """
    Main dispatcher – decides which sub-parser to call.
    """
    if not os.path.exists(file_path):
        raise ValueError("File does not exist")
    if os.path.getsize(file_path) == 0:
        raise ValueError("File is empty")

    ext = filename.lower().rsplit(".", 1)[-1]
    if ext in {"txt", "md", "json", "xml"}:
        return parse_txt(file_path)
    if ext == "docx":
        return parse_docx(file_path)
    if ext == "pptx":
        return parse_pptx(file_path)
    if ext == "pdf":
        return parse_pdf(file_path)
    if ext in {"csv", "xls", "xlsx"}:
        return parse_csv(file_path, filename)
    raise ValueError("File type not supported")


# ==========================================================
# PDF (text + images via PyMuPDF)
# ==========================================================
def parse_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    out = []

    for page_idx, page in enumerate(doc, start=1):
        out.append(f"--- Page {page_idx} ---\n")
        out.append(page.get_text())

        for img_idx, img in enumerate(page.get_images(full=True), start=1):
            xref = img[0]
            base = doc.extract_image(xref)
            img_bytes = base["image"]
            ocr_txt = ocr_image_bytes(img_bytes, f"PDF page {page_idx} image {img_idx}")
            if ocr_txt:
                out.append(f"\n[OCR – Page {page_idx} Image {img_idx}]\n{ocr_txt}\n")

    joined = "\n".join(out).strip()
    return joined if joined else "The file returned no content"


# ==========================================================
# DOCX (text + tables + images)
# ==========================================================
def parse_docx(file_path: str) -> str:
    document = Document(file_path)
    out = []

    # 1. Paragraphs & tables --------------------------------
    for item in document.iter_inner_content():
        if isinstance(item, Paragraph):
            out.append(item.text)
        elif isinstance(item, Table):
            for row in item.rows:
                cells = "\t".join(cell.text for cell in row.cells)
                out.append(cells)
        # newline after every element
        out.append("")

    # 2. Inline / header images ------------------------------
    for rel in document.part._rels.values():
        if "image" in rel.reltype:
            img_bytes = rel.target_part.blob
            ocr_txt = ocr_image_bytes(img_bytes, "DOCX image")
            if ocr_txt:
                out.append("[OCR – Image]\n" + ocr_txt + "\n")

    return "\n".join(out).strip()


# ==========================================================
# PPTX (text frames, tables, images)
# ==========================================================
def _collect_shape_text(shape, chunk_list):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            chunk_list.append("".join(run.text for run in para.runs))
    elif shape.has_table:
        for row in shape.table.rows:
            row_text = "\t".join(cell.text for cell in row.cells)
            chunk_list.append(row_text)


def _iter_shapes_recursive(shapes):
    for shp in shapes:
        yield shp
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_recursive(shp.shapes)


def parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    out = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_chunks = [f"--- Slide {slide_idx} ---"]
        for shp in _iter_shapes_recursive(slide.shapes):
            # regular text / table
            _collect_shape_text(shp, slide_chunks)

            # images
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                ocr_txt = ocr_image_bytes(shp.image.blob, f"PPTX slide {slide_idx} image")
                if ocr_txt:
                    slide_chunks.append(f"[OCR – Slide {slide_idx} Image]\n{ocr_txt}")

        out.append("\n".join(slide_chunks))
    return "\n\n".join(out).strip()


# ==========================================================
# Simple ASCII files
# ==========================================================
def parse_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


# ==========================================================
# CSV / Excel (unchanged)
# ==========================================================
def parse_csv(file_path: str, filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    with open(file_path, "rb") as f:
        data = f.read()

    if ext == ".csv":
        df = pd.read_csv(file_path)
        df_dict = None
    elif ext == ".xlsx":
        df_dict = pd.read_excel(io.BytesIO(data), sheet_name=None)
    elif ext == ".xls":
        wb = xlrd.open_workbook(file_contents=data, formatting_info=False)
        df_dict = {
            sheet.name: pd.DataFrame(
                sheet.get_rows(), columns=None
            )
            for sheet in wb.sheets()
        }
    else:
        raise ValueError(f"Unsupported extension: {ext}")

    if df_dict is None:
        body = df.to_json(orient="records", lines=True)
    else:
        parts = []
        for name, sheet_df in df_dict.items():
            parts.append(f"--- Sheet: {name} ---")
            parts.append(sheet_df.to_json(orient="records", lines=True))
        body = "\n".join(parts)

    preamble = (
        "Below is Excel data in the form of JSON, broken down by tabs. "
        "Depending on the ask, you may need to query the data. Ensure that "
        "all your calculations are correct, showing your thought process when applicable."
    )
    return f"{preamble}\n{body}"


# ==========================================================
# CLI helper
# ==========================================================
if __name__ == "__main__":
    if len(sys.argv) != 3:
        sys.stderr.write("Usage: parser_multi.py <file_path> <file_name>\n")
        sys.exit(1)

    print(parse_doc(sys.argv[1], sys.argv[2]))

